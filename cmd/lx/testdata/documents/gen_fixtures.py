#!/usr/bin/env python3

import io
import zipfile
import pathlib

out = pathlib.Path(__file__).parent

from fpdf import FPDF

pdf = FPDF()
pdf.add_page()
pdf.set_font("Helvetica", size=12)
pdf.cell(0, 10, "Hello PDF World")
pdf.ln()
pdf.cell(0, 10, "Second line of text")
pdf.output(out / "sample.pdf")
print("wrote sample.pdf")

from docx import Document

doc = Document()
doc.add_paragraph("Hello DOCX World")
doc.add_paragraph("Second paragraph")
doc.save(out / "sample.docx")
print("wrote sample.docx")

from openpyxl import Workbook

wb = Workbook()
ws = wb.active
ws.title = "Sheet1"
ws["A1"] = "Name"
ws["B1"] = "Score"
ws["A2"] = "Alice"
ws["B2"] = 42
wb.save(out / "sample.xlsx")
print("wrote sample.xlsx")

from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Hello PPTX World"
slide.placeholders[1].text = "First slide content"
slide2 = prs.slides.add_slide(slide_layout)
slide2.shapes.title.text = "Second Slide"
slide2.placeholders[1].text = "Second slide content"
prs.save(out / "sample.pptx")
print("wrote sample.pptx")

def write_odf_zip(path, mimetype, content_xml):
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        info = zipfile.ZipInfo("mimetype")
        info.compress_type = zipfile.ZIP_STORED
        zf.writestr(info, mimetype)
        manifest_xml = f"""\
<?xml version="1.0" encoding="UTF-8"?>
<manifest:manifest xmlns:manifest="urn:oasis:names:tc:opendocument:xmlns:manifest:1.0">
  <manifest:file-entry manifest:full-path="/" manifest:media-type="{mimetype}"/>
  <manifest:file-entry manifest:full-path="content.xml" manifest:media-type="text/xml"/>
</manifest:manifest>
"""
        zf.writestr("META-INF/manifest.xml", manifest_xml)
        zf.writestr("content.xml", content_xml)
    path.write_bytes(buf.getvalue())

odt_content_xml = """\
<?xml version="1.0" encoding="UTF-8"?>
<office:document-content
    xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"
    xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">
  <office:body>
    <office:text>
      <text:p>Hello ODT World</text:p>
      <text:p>Second paragraph</text:p>
    </office:text>
  </office:body>
</office:document-content>
"""
write_odf_zip(out / "sample.odt", "application/vnd.oasis.opendocument.text", odt_content_xml)
print("wrote sample.odt")

ods_content_xml = """\
<?xml version="1.0" encoding="UTF-8"?>
<office:document-content
    xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"
    xmlns:table="urn:oasis:names:tc:opendocument:xmlns:table:1.0"
    xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">
  <office:body>
    <office:spreadsheet>
      <table:table table:name="Sheet1">
        <table:table-row>
          <table:table-cell><text:p>Hello ODS World</text:p></table:table-cell>
          <table:table-cell><text:p>Value</text:p></table:table-cell>
        </table:table-row>
      </table:table>
    </office:spreadsheet>
  </office:body>
</office:document-content>
"""
write_odf_zip(out / "sample.ods", "application/vnd.oasis.opendocument.spreadsheet", ods_content_xml)
print("wrote sample.ods")

odp_content_xml = """\
<?xml version="1.0" encoding="UTF-8"?>
<office:document-content
    xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"
    xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    xmlns:draw="urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
    xmlns:presentation="urn:oasis:names:tc:opendocument:xmlns:presentation:1.0">
  <office:body>
    <office:presentation>
      <draw:page draw:name="page1">
        <draw:text-box>
          <text:p>Hello ODP World</text:p>
          <text:p>Second paragraph</text:p>
        </draw:text-box>
      </draw:page>
    </office:presentation>
  </office:body>
</office:document-content>
"""
write_odf_zip(out / "sample.odp", "application/vnd.oasis.opendocument.presentation", odp_content_xml)
print("wrote sample.odp")
